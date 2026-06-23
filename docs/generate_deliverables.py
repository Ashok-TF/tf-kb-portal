"""Generate Word implementation doc and PowerPoint deck for Sonu presentation."""

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.util import Inches as PptInches, Pt as PptPt

OUT_DIR = Path(__file__).parent


def add_heading(doc, text, level=1):
    return doc.add_heading(text, level=level)


def add_para(doc, text, bold=False):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = bold
    return p


def add_bullet(doc, text, level=0):
    return doc.add_paragraph(text, style="List Bullet")


def build_docx():
    doc = Document()

    # Title page
    title = doc.add_heading("Enterprise Knowledge Base Platform", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub = doc.add_paragraph("v1 Implementation Document")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.runs[0].bold = True
    sub.runs[0].font.size = Pt(16)

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(
        "Prepared for: Sonu\n"
        "Team: Ashok (Frontend) + Sreepad (Backend)\n"
        "Timeline: ~6 weeks to v1 demo\n"
        "Reference: kb-platform-animation.html (6-step architecture walkthrough)"
    )

    doc.add_page_break()

    # 1. Executive Summary
    add_heading(doc, "1. Executive Summary", 1)
    doc.add_paragraph(
        "We are building an Enterprise Knowledge Base Platform based on the architecture "
        "defined in the HTML walkthrough. The platform lets employees upload knowledge, "
        "organize it by domain (HR, Engineering, Sales, etc.), and ask questions in natural "
        "language — with AI-generated answers backed by source citations."
    )
    doc.add_paragraph(
        "v1 (6 weeks) delivers the core intelligence layer: ingestion of documents, semantic "
        "storage, and agentic search with grounded answers (HTML Step 05)."
    )
    doc.add_paragraph(
        "v2 and v3 add enrichment, connectors, knowledge graph, and enterprise governance "
        "(HTML Steps 02, 03, 04, 06)."
    )

    # 2. Reference Architecture
    add_heading(doc, "2. Reference Architecture (from HTML Walkthrough)", 1)
    doc.add_paragraph(
        "The HTML file (kb-platform-animation.html) is a 6-scene animated architecture spec. "
        "Each scene represents a pipeline stage:"
    )

    table = doc.add_table(rows=7, cols=3)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    hdr[0].text = "Step"
    hdr[1].text = "HTML Title"
    hdr[2].text = "What It Depicts"
    rows = [
        ("01", "Ingestion — Capture knowledge from anywhere",
         "Documents (PDF/DOCX/PPTX/XLSX/IMG), Website crawling, Artifacts & URLs → Content ingestion"),
        ("02", "Enrichment — Understand before you index",
         "OCR, Vision-LLM, Entity extraction, Relationships + summaries"),
        ("03", "Knowledge artifacts — Turn content into reusable knowledge",
         "Vector index, Wiki articles, Entities, Relationships, Knowledge graph → KB registry"),
        ("04", "Search framework — Pluggable search, swappable storage",
         "Pinecone, OpenSearch, Qdrant, PGVector, Azure AI Search → Multiple search strategies"),
        ("05", "Agentic search — The agent reasons over merged knowledge",
         "User query → AI agent → parallel multi-source search → fusion → grounded answer with sources"),
        ("06", "Governance & growth — Secure, audited, built to extend",
         "KB isolation, audit logging, usage analytics, connector framework (SharePoint, Confluence, etc.)"),
    ]
    for i, (step, title, desc) in enumerate(rows, 1):
        table.rows[i].cells[0].text = step
        table.rows[i].cells[1].text = title
        table.rows[i].cells[2].text = desc

    doc.add_paragraph()
    add_para(doc, "Tip for presentation: Open kb-platform-animation.html in a browser and walk through Steps 1–6 live.", bold=True)

    # 3. Version Roadmap
    add_heading(doc, "3. Version Roadmap — What We Achieve When", 1)

    table2 = doc.add_table(rows=7, cols=4)
    table2.style = "Table Grid"
    h2 = table2.rows[0].cells
    h2[0].text = "HTML Step"
    h2[1].text = "v1 (~6 weeks)"
    h2[2].text = "v2 (~+8 weeks)"
    h2[3].text = "v3 (~+10 weeks)"
    vrows = [
        ("01 Ingestion", "Documents upload only", "URLs + website crawler", "Full sources + artifacts"),
        ("02 Enrichment", "Basic text extraction + OCR", "LLM summaries, entities", "Vision-LLM for diagrams"),
        ("03 Knowledge artifacts", "Vector index + KB registry", "Wiki articles, entities", "Knowledge graph"),
        ("04 Search framework", "pgvector semantic search", "Hybrid search, pluggable stores", "Graph/wiki/context search"),
        ("05 Agentic search", "FULL v1 FOCUS — agent + RAG + citations", "Multi-source parallel retrieval", "Advanced fusion & re-ranking"),
        ("06 Governance", "Basic per-KB isolation", "Audit logging + SharePoint", "Analytics + Confluence/GitHub"),
    ]
    for i, row in enumerate(vrows, 1):
        for j, val in enumerate(row):
            table2.rows[i].cells[j].text = val

    # 4. v1 Detailed
    add_heading(doc, "4. v1 — What We Are Building", 1)

    add_heading(doc, "4.1 v1 Goal", 2)
    doc.add_paragraph(
        "In 6 weeks, deliver a working KB platform where users upload documents, organize them "
        "into Knowledge Bases, ask questions in plain English, and receive AI answers with "
        "verifiable source citations — covering HTML Steps 01 (partial), 03 (partial), 04 (partial), "
        "and 05 (core)."
    )

    add_heading(doc, "4.2 v1 User Journey", 2)
    for step in [
        "LOGIN — User signs in to the portal",
        "CREATE KB — e.g. 'Pre-sales KB', 'HR Policies'",
        "UPLOAD DOCS — PDF, DOCX, PPTX, XLSX, images",
        "PROCESSING — Extract text → chunk → embed → store (automatic)",
        "ASK QUESTION — 'Where is the latest AI sales deck?'",
        "AGENT WORKS — Picks relevant KB(s) → searches vectors → fuses results",
        "GET ANSWER — Grounded response + citations (filename, excerpt, score)",
    ]:
        add_bullet(doc, step)

    add_heading(doc, "4.3 v1 Architecture", 2)
    doc.add_paragraph(
        "Browser → Next.js Frontend → FastAPI Backend → PostgreSQL (metadata) + pgvector "
        "(embeddings/search) + OpenAI (embeddings + chat LLM)"
    )

    add_heading(doc, "4.4 v1 Features", 2)
    feat_table = doc.add_table(rows=15, cols=3)
    feat_table.style = "Table Grid"
    feat_table.rows[0].cells[0].text = "#"
    feat_table.rows[0].cells[1].text = "Feature"
    feat_table.rows[0].cells[2].text = "HTML Step"
    features = [
        ("F1", "User authentication", "—"),
        ("F2", "Knowledge Base CRUD", "03"),
        ("F3", "Document upload (PDF, DOCX, PPTX, XLSX, images)", "01"),
        ("F4", "Ingestion pipeline (extract → chunk → embed → store)", "01"),
        ("F5", "Processing status (pending → processing → ready/failed)", "01"),
        ("F6", "Vector storage (pgvector in Postgres)", "03, 04"),
        ("F7", "Semantic search", "04"),
        ("F8", "KB registry agent (auto-select relevant KBs)", "05"),
        ("F9", "RAG retrieval (top-K chunks)", "05"),
        ("F10", "Answer fusion (dedup, re-rank, compress)", "05"),
        ("F11", "Grounded LLM answer", "05"),
        ("F12", "Citations (filename, excerpt, score)", "05"),
        ("F13", "Chat UI (per-KB + global)", "05"),
        ("F14", "Per-KB isolation", "06"),
    ]
    for i, (fid, feat, step) in enumerate(features, 1):
        feat_table.rows[i].cells[0].text = fid
        feat_table.rows[i].cells[1].text = feat
        feat_table.rows[i].cells[2].text = step

    add_heading(doc, "4.5 v1 — NOT Included", 2)
    for item in [
        "Website crawling, URL ingestion",
        "Vision-LLM, entity/relationship extraction",
        "Wiki auto-generation, knowledge graph",
        "Multiple vector DB providers (Pinecone, OpenSearch, etc.)",
        "SharePoint / Confluence / GitHub connectors",
        "Audit logs, usage analytics, SSO",
    ]:
        add_bullet(doc, item)

    # 5. Task Breakdown
    add_heading(doc, "5. Complete Task Breakdown", 1)

    phases = [
        ("Phase 0 — Project Setup (Week 1)", [
            "0.1 Repo setup, branch strategy, env templates (Both, 0.5d)",
            "0.2 Architecture doc + API contract (Both, 1d)",
            "0.3 Docker Compose: Postgres + pgvector + backend + frontend (Sreepad, 2d)",
            "0.4 Postgres schema: users, KBs, documents, chunks (Sreepad, 2d)",
            "0.5 pgvector extension + vector table design (Sreepad, 1d)",
            "0.6 OpenAI embedder integration (Sreepad, 1d)",
            "0.7 OpenAI chat/LLM provider abstraction (Sreepad, 1d)",
            "0.8 Frontend scaffold: Next.js, layout, auth shell (Ashok, 2d)",
            "0.9 Design system: colors, components (Ashok, 2d)",
        ]),
        ("Phase 1 — Ingestion (Week 1–2) → HTML Step 01", [
            "1.1 Auth API: login, token, multi-user (Sreepad, 2d)",
            "1.2 KB CRUD API (Sreepad, 1d)",
            "1.3 Document upload API — multipart, 202 async (Sreepad, 2d)",
            "1.4 Text extraction: PDF, DOCX, PPTX, XLSX, TXT, MD (Sreepad, 2d)",
            "1.5 Image OCR — Tesseract (Sreepad, 1d)",
            "1.6 Chunking service — ~900 chars, 150 overlap (Sreepad, 1d)",
            "1.7 Embed + upsert to pgvector (Sreepad, 2d)",
            "1.8 Document status API + background job queue (Sreepad, 2d)",
            "1.9 Login page (Ashok, 1d)",
            "1.10 KB dashboard: list, create, delete (Ashok, 2d)",
            "1.11 KB detail page + document table (Ashok, 2d)",
            "1.12 Upload zone: drag-drop, progress, status poll (Ashok, 2d)",
        ]),
        ("Phase 2 — Search & RAG (Week 2–3) → HTML Steps 04 + 05", [
            "2.1 Semantic search API: query → embed → pgvector top-K (Sreepad, 2d)",
            "2.2 RAG prompt template + citation extraction (Sreepad, 2d)",
            "2.3 Chat API: single-KB POST /kb/{id}/chat (Sreepad, 2d)",
            "2.4 Search results UI — snippet cards (Ashok, 1d)",
            "2.5 Chat UI: messages, input, loading (Ashok, 3d)",
            "2.6 Citation cards in chat (Ashok, 2d)",
        ]),
        ("Phase 3 — Agentic Search (Week 3–4) → HTML Step 05 Core", [
            "3.1 KB scoring agent: match query to KB name/description (Sreepad, 2d)",
            "3.2 Multi-KB parallel retrieval (Sreepad, 2d)",
            "3.3 Fusion layer: dedup, re-rank, context window (Sreepad, 2d)",
            "3.4 Global chat API: POST /chat (Sreepad, 1d)",
            "3.5 Global chat page — auto KB selection (Ashok, 2d)",
            "3.6 KB badge on citations (Ashok, 1d)",
            "3.7 Empty/error states (Ashok, 1d)",
        ]),
        ("Phase 4 — Polish & Demo (Week 4–5)", [
            "4.1 RAG quality tuning: prompt, top_k, chunk size (Sreepad, 2d)",
            "4.2 Reindex + delete document APIs (Sreepad, 1d)",
            "4.3 Reindex/delete UI actions (Ashok, 1d)",
            "4.4 App header, nav, responsive layout (Ashok, 2d)",
            "4.5 Embed HTML walkthrough as /demo page (Ashok, 1d)",
            "4.6 Sample KBs + test documents for demo (Both, 1d)",
            "4.7 End-to-end smoke test + bug fixes (Both, 2d)",
            "4.8 Setup README + runbook (Both, 1d)",
        ]),
        ("Phase 5 — Buffer & Cloud Prep (Week 6)", [
            "5.1 Production Docker images (Sreepad, 2d)",
            "5.2 AWS deployment guide: RDS, ECS/EC2, S3 (Sreepad, 2d)",
            "5.3 UI fixes from stakeholder feedback (Ashok, 2d)",
            "5.4 v2 backlog document (Both, 1d)",
        ]),
    ]

    for phase_title, tasks in phases:
        add_heading(doc, phase_title, 2)
        for t in tasks:
            add_bullet(doc, t)

    # 6. v2 & v3
    add_heading(doc, "6. v2 & v3 — What Comes Next", 1)

    add_heading(doc, "v2 (~8 weeks after v1)", 2)
    for item in [
        "Enrichment (Step 02): LLM summaries, entity extraction, relationship tagging",
        "Ingestion expansion (Step 01): URL ingestion, website crawler",
        "Wiki articles (Step 03): Auto-generated from documents",
        "Search upgrade (Step 04): Hybrid search (vector + keyword BM25)",
        "Governance start (Step 06): Audit log (upload, search, access events)",
        "First connector (Step 06): SharePoint (Azure AD OAuth, document sync)",
    ]:
        add_bullet(doc, item)

    add_heading(doc, "v3 (~10 weeks after v2)", 2)
    for item in [
        "Knowledge graph (Step 03): Entities + relationships in Postgres",
        "Advanced search (Step 04): Graph, wiki, context search",
        "Full agentic (Step 05): Parallel multi-source retrieval + advanced fusion",
        "Vision-LLM (Step 02): Diagram/table/chart understanding",
        "Connectors (Step 06): Confluence, GitHub, Jira",
        "Enterprise (Step 06): Usage analytics dashboard, SSO/Azure AD",
    ]:
        add_bullet(doc, item)

    add_heading(doc, "Full Vision Timeline", 2)
    tl = doc.add_table(rows=4, cols=3)
    tl.style = "Table Grid"
    tl.rows[0].cells[0].text = "Milestone"
    tl.rows[0].cells[1].text = "Duration"
    tl.rows[0].cells[2].text = "Cumulative"
    tl_data = [
        ("v1 — Agentic chat + upload", "6 weeks", "6 weeks"),
        ("v2 — Enrichment + SharePoint", "8 weeks", "~3.5 months"),
        ("v3 — Graph + full enterprise", "10 weeks", "~6 months"),
    ]
    for i, row in enumerate(tl_data, 1):
        for j, val in enumerate(row):
            tl.rows[i].cells[j].text = val

    # 7. Team & Tech
    add_heading(doc, "7. Team, Tech Stack & Prerequisites", 1)

    add_heading(doc, "Team", 2)
    team = doc.add_table(rows=3, cols=3)
    team.style = "Table Grid"
    team.rows[0].cells[0].text = "Person"
    team.rows[0].cells[1].text = "Role"
    team.rows[0].cells[2].text = "v1 Ownership"
    team.rows[1].cells[0].text = "Ashok"
    team.rows[1].cells[1].text = "Frontend Lead"
    team.rows[1].cells[2].text = "All UI: login, KB dashboard, upload, chat, citations, demo page"
    team.rows[2].cells[0].text = "Sreepad"
    team.rows[2].cells[1].text = "Backend Lead"
    team.rows[2].cells[2].text = "All API: auth, ingestion, pgvector, agent, RAG, Docker"

    add_heading(doc, "Tech Stack (v1)", 2)
    tech = doc.add_table(rows=7, cols=2)
    tech.style = "Table Grid"
    tech.rows[0].cells[0].text = "Layer"
    tech.rows[0].cells[1].text = "Technology"
    tech_data = [
        ("Frontend", "Next.js, React, TypeScript, Tailwind CSS"),
        ("Backend", "Python, FastAPI, SQLAlchemy"),
        ("Database", "PostgreSQL + pgvector"),
        ("AI", "OpenAI (text-embedding-3-small + gpt-4o-mini)"),
        ("Queue", "Background tasks (Redis/RQ optional)"),
        ("DevOps", "Docker Compose (local), AWS (staging post-v1)"),
    ]
    for i, row in enumerate(tech_data, 1):
        tech.rows[i].cells[0].text = row[0]
        tech.rows[i].cells[1].text = row[1]

    add_heading(doc, "Prerequisites Before Kickoff", 2)
    for item in [
        "OpenAI API key with billing enabled",
        "GitHub repo created",
        "Postgres with pgvector (via Docker)",
        "6-week allocation for Ashok + Sreepad",
        "Stakeholder demo slot (Week 5–6)",
    ]:
        add_bullet(doc, item)

    # 8. Success Criteria
    add_heading(doc, "8. v1 Success Criteria", 1)
    for i, criterion in enumerate([
        "User logs in securely",
        "Creates 2+ Knowledge Bases",
        "Uploads PDF/DOCX/PPTX successfully",
        "Documents reach 'ready' status",
        "Asks natural-language question in chat",
        "Agent selects correct KB automatically",
        "Answer is relevant and grounded in documents",
        "2+ citations shown with correct filenames",
        "Runs locally via Docker Compose",
        "HTML walkthrough demo page available at /demo",
    ], 1):
        add_bullet(doc, f"{i}. {criterion}")

    # 9. Risks
    add_heading(doc, "9. Risks & Mitigations", 1)
    risk = doc.add_table(rows=6, cols=3)
    risk.style = "Table Grid"
    risk.rows[0].cells[0].text = "Risk"
    risk.rows[0].cells[1].text = "Impact"
    risk.rows[0].cells[2].text = "Mitigation"
    risk_data = [
        ("OpenAI API costs", "Medium", "Use gpt-4o-mini; monitor usage"),
        ("RAG answer quality", "High", "Tune prompts, chunk size, top_k in Week 5"),
        ("SharePoint complexity (v2)", "High", "Defer to v2; manual upload in v1"),
        ("Scope creep", "High", "Strict v1 sign-off; v2 backlog for extras"),
        ("OCR quality on images", "Medium", "Clear error messages; optional in v1"),
    ]
    for i, row in enumerate(risk_data, 1):
        for j, val in enumerate(row):
            risk.rows[i].cells[j].text = val

    # 10. One-liner
    add_heading(doc, "10. One-Line Pitch", 1)
    p = doc.add_paragraph()
    run = p.add_run(
        '"In 6 weeks we deliver an Enterprise KB where employees ask questions in plain English '
        'and get AI answers backed by source documents — the foundation for SharePoint, knowledge '
        'graph, and full enterprise features."'
    )
    run.italic = True

    out = OUT_DIR / "Enterprise-KB-Platform-v1-Implementation.docx"
    doc.save(out)
    return out


def set_slide_title(slide, title, subtitle=None):
    title_shape = slide.shapes.title
    title_shape.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def add_content_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = bullet
        p.level = 0
    return slide


def build_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    # Slide 1 — Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Enterprise Knowledge Base Platform"
    slide.placeholders[1].text = (
        "Implementation Plan v1\n\n"
        "Ashok + Sreepad · 6 weeks · ThoughtFocus\n"
        "Reference: kb-platform-animation.html (6-step architecture)"
    )

    slides_data = [
        ("The Vision — 6-Step Architecture", [
            "Step 01: Ingestion — capture from anywhere",
            "Step 02: Enrichment — AI understands content",
            "Step 03: Knowledge artifacts — vectors, wiki, graph",
            "Step 04: Search framework — pluggable retrieval",
            "Step 05: Agentic search — AI answers with sources ★",
            "Step 06: Governance — security, audit, connectors",
            "",
            "Demo the HTML file live: kb-platform-animation.html",
        ]),
        ("The Problem", [
            "Company knowledge is scattered across PDFs, decks, policies",
            "Keyword search is not enough",
            "Employees need instant, trusted answers with proof",
        ]),
        ("The Solution", [
            "Upload → Organize → Ask → Get cited AI answers",
            "",
            'Example: "Where is the latest AI sales deck?"',
            "→ System finds the right Knowledge Base",
            "→ Searches documents",
            "→ Returns answer with source files",
        ]),
        ("Version Plan", [
            "v1: Upload + Agentic chat + citations — 6 weeks",
            "v2: Enrichment + SharePoint + audit — +8 weeks",
            "v3: Knowledge graph + full enterprise — +10 weeks",
            "",
            "Full HTML vision: ~6 months with 2 developers",
        ]),
        ("What v1 Covers (HTML Mapping)", [
            "Step 01 Ingestion      — Documents upload (partial)",
            "Step 02 Enrichment     — Basic text + OCR only (minimal)",
            "Step 03 Artifacts      — Vector index + KB registry (partial)",
            "Step 04 Search         — pgvector semantic search (partial)",
            "Step 05 Agentic search — ★ CORE v1 FOCUS ★",
            "Step 06 Governance     — Per-KB isolation only (minimal)",
        ]),
        ("v1 User Journey", [
            "1. Login",
            "2. Create Knowledge Base (e.g. Pre-sales, HR)",
            "3. Upload documents (PDF, DOCX, PPTX, XLSX, images)",
            "4. Auto-process: extract → chunk → embed → store",
            "5. Ask question in chat",
            "6. Agent picks relevant KB(s)",
            "7. Get AI answer with source citations",
        ]),
        ("v1 Architecture", [
            "Browser → Next.js UI → FastAPI API",
            "                          ↓",
            "              Postgres + pgvector + OpenAI",
            "",
            "Postgres — metadata (KBs, documents)",
            "pgvector — semantic search",
            "OpenAI — embeddings + chat answers",
        ]),
        ("6-Week Timeline", [
            "Week 1: Infrastructure + DB + OpenAI wired up",
            "Week 2: Upload pipeline + single-KB chat with citations",
            "Week 3: Multi-KB agent + global chat",
            "Week 4: UI polish + HTML walkthrough demo page",
            "Week 5: Testing + sample KBs + internal demo",
            "Week 6: AWS prep + v2 planning",
        ]),
        ("Team", [
            "Ashok (Frontend):",
            "  Login, KB dashboard, upload UI, chat UI, citations, design",
            "",
            "Sreepad (Backend):",
            "  Auth, ingestion pipeline, pgvector, RAG, agent, Docker",
        ]),
        ("v1 Success = Demo", [
            "✓ 2 Knowledge Bases with real documents",
            "✓ Natural-language question in chat",
            "✓ AI answer with 2+ source citations",
            "✓ Correct files referenced",
            "✓ Runs locally via Docker Compose",
        ]),
        ("v2 & v3 Preview", [
            "v2 (+8 weeks):",
            "  SharePoint connector · LLM enrichment · wiki · audit logging",
            "",
            "v3 (+10 weeks):",
            "  Knowledge graph · Vision-LLM · Confluence/GitHub · analytics",
        ]),
        ("What We Need", [
            "1. Approve v1 scope (agentic chat, 6 weeks)",
            "2. OpenAI API budget for dev + demo",
            "3. Ashok + Sreepad focused for 6 weeks",
            "4. Stakeholder demo slot — Week 5–6",
        ]),
        ("Closing", [
            '"In 6 weeks we deliver an Enterprise KB where employees ask questions',
            'in plain English and get AI answers backed by source documents —',
            "the foundation for SharePoint, knowledge graph, and full enterprise features.",
            "",
            "Questions?",
        ]),
    ]

    for title, bullets in slides_data:
        add_content_slide(prs, title, bullets)

    out = OUT_DIR / "Enterprise-KB-Platform-v1-Sonu-Deck.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    docx_path = build_docx()
    pptx_path = build_pptx()
    print(f"Created: {docx_path}")
    print(f"Created: {pptx_path}")
