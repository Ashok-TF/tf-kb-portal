"""Text extraction for a variety of file types.

Each extractor is best-effort and degrades gracefully: if an optional library
or system dependency (e.g. the Tesseract OCR binary) is missing, we raise a
clear error that gets stored on the document instead of crashing the worker.
"""

from __future__ import annotations

import csv
from pathlib import Path


class ExtractionError(Exception):
    pass


def _ext(path: Path) -> str:
    return path.suffix.lower().lstrip(".")


def extract_text(path: Path, content_type: str | None = None) -> str:
    """Return plain text extracted from the file at `path`."""
    ext = _ext(path)

    if ext == "pdf":
        return _extract_pdf(path)
    if ext in {"txt", "md", "markdown", "log", "json"}:
        return _read_text(path)
    if ext == "csv":
        return _extract_csv(path)
    if ext in {"docx"}:
        return _extract_docx(path)
    if ext in {"pptx"}:
        return _extract_pptx(path)
    if ext in {"xlsx", "xlsm"}:
        return _extract_xlsx(path)
    if ext in {"xls"}:
        # Legacy Excel format isn't supported by openpyxl.
        raise ExtractionError(
            "Legacy .xls files are not supported. Please convert to .xlsx and re-upload."
        )
    if ext in {"png", "jpg", "jpeg", "gif", "bmp", "tiff", "webp"}:
        return _extract_image(path)

    # Unknown type: try to read as UTF-8 text, else fail clearly.
    try:
        return _read_text(path)
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(f"Unsupported file type: .{ext}") from exc


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_csv(path: Path) -> str:
    rows: list[str] = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
        for row in csv.reader(f):
            rows.append(" | ".join(cell.strip() for cell in row))
    return "\n".join(rows)


def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError("pypdf is not installed") from exc

    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text)
    return "\n\n".join(parts)


def _extract_docx(path: Path) -> str:
    try:
        import docx
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError("python-docx is not installed") from exc

    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs if p.text and p.text.strip()]

    # Also pull text out of tables.
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError("python-pptx is not installed") from exc

    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        slide_parts.append(line)
        if slide_parts:
            parts.append(f"[Slide {idx}]\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError("openpyxl is not installed") from exc

    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def _extract_image(path: Path) -> str:
    """OCR an image. Requires the Tesseract binary + pytesseract."""
    try:
        from PIL import Image
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError("Pillow is not installed") from exc

    try:
        import pytesseract
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError("pytesseract is not installed") from exc

    try:
        with Image.open(path) as img:
            text = pytesseract.image_to_string(img)
    except pytesseract.TesseractNotFoundError as exc:
        raise ExtractionError(
            "Tesseract OCR engine not found. Install it (e.g. `choco install tesseract` "
            "on Windows or `apt-get install tesseract-ocr` on Linux) to extract text from images."
        ) from exc
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(f"Failed to OCR image: {exc}") from exc

    if not text.strip():
        raise ExtractionError("No text could be extracted from the image via OCR.")
    return text
